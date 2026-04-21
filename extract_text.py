#!/usr/bin/env python3
"""Extract text from all unique attachments so we can discuss the content."""
import os
from docx import Document
from pptx import Presentation

ATT = "/home/amsamms/projects/EMBA/marketing/midterm_preparation/attachments"
OUT = "/home/amsamms/projects/EMBA/marketing/midterm_preparation/extracted"
os.makedirs(OUT, exist_ok=True)

# Use only the first (EMBA_marketing_ prefix) copies for docx since email 1 and 2 are dupes
files = sorted(os.listdir(ATT))
seen_base = set()
unique = []
for f in files:
    base = f.split("__", 1)[1] if "__" in f else f
    if base in seen_base:
        continue
    seen_base.add(base)
    unique.append(f)

for f in unique:
    path = os.path.join(ATT, f)
    base = f.split("__", 1)[1] if "__" in f else f
    out_path = os.path.join(OUT, base + ".txt")
    print(f"\n##### {base} #####")
    lines = []
    if f.lower().endswith(".docx"):
        doc = Document(path)
        for p in doc.paragraphs:
            if p.text.strip():
                lines.append(p.text)
        for t in doc.tables:
            for row in t.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    lines.append(row_text)
    elif f.lower().endswith(".pptx") or f.lower().endswith(".ppt"):
        try:
            prs = Presentation(path)
        except Exception as e:
            print(f"  cannot open (likely legacy .ppt): {e}")
            continue
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in para.runs).strip()
                        if text:
                            lines.append(text)
    content = "\n".join(lines)
    with open(out_path, "w") as g:
        g.write(content)
    print(content[:400])
    print(f"... ({len(content)} chars total -> {out_path})")
